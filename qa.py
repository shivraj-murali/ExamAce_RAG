from flask import Flask, request, jsonify
from langchain_community.document_loaders.pdf import PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_mongodb import MongoDBAtlasVectorSearch
from pymongo import MongoClient
from langchain_huggingface import HuggingFaceEndpoint
from dotenv import load_dotenv
from flask_cors import CORS, cross_origin
from pptx import Presentation
from pptx.util import Inches
import google.generativeai as genai
import os

load_dotenv()

embeddings = HuggingFaceEmbeddings(model_name="mixedbread-ai/mxbai-embed-large-v1", encode_kwargs={'precision': 'binary'})

# Vector Store
one_bit_vectorstore = FAISS.load_local("RL", embeddings, allow_dangerous_deserialization=True)
retriever = one_bit_vectorstore.as_retriever(search_kwargs={"k": 10})  # Increased context retrieval


one_bit_vectorstore_dwm = FAISS.load_local("DWM-dataStore", embeddings, allow_dangerous_deserialization=True)
retriever_dwm = one_bit_vectorstore_dwm.as_retriever(search_kwargs={"k": 10})  # Increased context retrieval

app = Flask(__name__)
cors = CORS(app)
app.config['CORS_HEADERS'] = 'Content-Type'

# Configure the Gemini API
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=GOOGLE_API_KEY)

def create_slide(presentation, title, content_points):
    slide_layout = presentation.slide_layouts[1]  # Use the layout with title and content
    slide = presentation.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    
    for point in content_points:
        p = content_placeholder.text_frame.add_paragraph()
        p.text = point
        p.level = 0  # Bullet point level

def generate_presentation(slide_data, template_path, output_path):
    presentation = Presentation(template_path)
    
    for slide_info in slide_data:
        title = slide_info['title']
        content_points = slide_info['content']
        
        while content_points:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title_placeholder = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            
            # Add points to the slide until it is full
            remaining_points = []
            for point in content_points:
                p = content_placeholder.text_frame.add_paragraph()
                p.text = point
                p.level = 0  # Bullet point level
                
                # Check if the content exceeds the slide's capacity
                if content_placeholder.text_frame.fit_text():
                    remaining_points.append(point)
                    content_placeholder.text_frame.text = content_placeholder.text_frame.text.rsplit('\n', 1)[0]
                    break
            
            content_points = remaining_points
    
    presentation.save(output_path)

def is_comparative_or_exploratory_question(question):
    """
    Detect if the question is comparative or exploratory
    """
    comparative_keywords = [
        'difference', 'compare', 'contrast', 'best', 'better', 'vs', 'versus', 
        'what is', 'explain', 'describe', 'how', 'why', 'characteristics', 
        'features', 'advantages', 'disadvantages'
    ]
    
    # Convert question to lowercase for case-insensitive matching
    lower_question = question.lower()
    
    # Check if any comparative keywords are in the question
    return any(keyword in lower_question for keyword in comparative_keywords)

@app.route("/cn", methods=['POST'])
def hello_world():
    content = request.json

    # Gemini model configuration
    model = genai.GenerativeModel('gemini-1.5-pro')

    # Improved template with flexible answering strategy
    template = """
    Context: {context}
    
    Question: {question}
    
Guidelines for Response:

1. If context provides direct information about the question, use that information and cite the specific section or document name as your source.

2. If the question is comparative or exploratory, provide a comprehensive explanation with citations for each key point from specific sections of your knowledge base.

3. If no specific context is available, give a general, informative response based on broader knowledge, clearly indicating which portions are from general knowledge.

4. Be precise, clear, and provide structured information with proper attribution throughout.

5. If absolutely no information is available, clearly state that and explain the limitations of your knowledge on this topic.

6. Always end your response with a "Sources:" section that lists all documents, sections, or knowledge sources used to generate the answer.

7. For each claim or piece of information in your response, include a reference number [1], [2], etc. that corresponds to the source in your citation list.
    """

    # Add basic validation for question content
    if 'question' not in content or not content['question'].strip():
        return {"answer": "Please provide a valid question."}, 400

    # Add error handling for the API call
    try:
        # Retrieve context first
        context_docs = retriever.invoke(content['question'])
        context_text = "\n".join([doc.page_content for doc in context_docs])

        # Determine if it's a comparative or exploratory question
        is_comparative = is_comparative_or_exploratory_question(content['question'])

        # Prepare the full prompt
        full_prompt = template.format(context=context_text, question=content['question'])

        # If it's a comparative or exploratory question, add more context
        if is_comparative and context_text:
            full_prompt += "\n\nNote: This is a comparative or exploratory question. Provide a comprehensive analysis using available context and broader knowledge."

        # Generate response using Gemini
        response = model.generate_content(full_prompt)
        
        return {"answer": response.text}
    except Exception as e:
        return {"answer": f"An error occurred while processing your request. Please try again with a more specific question.", "error": str(e)}, 500


@app.route("/dwm", methods=['POST'])
def dwm():
    content = request.json

    # Gemini model configuration
    model = genai.GenerativeModel('gemini-1.5-pro')

    # Improved template with flexible answering strategy
    template = """
    Context: {context}
    
    Question: {question}
    
Guidelines for Response:

1. If context provides direct information about the question, use that information and cite the specific section or document name as your source.

2. If the question is comparative or exploratory, provide a comprehensive explanation with citations for each key point from specific sections of your knowledge base.

3. If no specific context is available, give a general, informative response based on broader knowledge, clearly indicating which portions are from general knowledge.

4. Be precise, clear, and provide structured information with proper attribution throughout.

5. If absolutely no information is available, clearly state that and explain the limitations of your knowledge on this topic.

6. Always end your response with a "Sources:" section that lists all documents, sections, or knowledge sources used to generate the answer.

7. For each claim or piece of information in your response, include a reference number [1], [2], etc. that corresponds to the source in your citation list.
    """

    # Add basic validation for question content
    if 'question' not in content or not content['question'].strip():
        return {"answer": "Please provide a valid question."}, 400

    # Add error handling for the API call
    try:
        # Retrieve context first
        context_docs = retriever_dwm.invoke(content['question'])
        context_text = "\n".join([doc.page_content for doc in context_docs])

        # Determine if it's a comparative or exploratory question
        is_comparative = is_comparative_or_exploratory_question(content['question'])

        # Prepare the full prompt
        full_prompt = template.format(context=context_text, question=content['question'])

        # If it's a comparative or exploratory question, add more context
        if is_comparative and context_text:
            full_prompt += "\n\nNote: This is a comparative or exploratory question. Provide a comprehensive analysis using available context and broader knowledge."

        # Generate response using Gemini
        response = model.generate_content(full_prompt)
        
        return {"answer": response.text}
    except Exception as e:
        return {"answer": f"An error occurred while processing your request. Please try again with a more specific question.", "error": str(e)}, 500




@app.route('/ppt', methods=['GET'])
def ppt():
    slides_data = [
        {
            "title": "Introduction to AI",
            "content": [
                "Definition of Artificial Intelligence",
                "History of AI development",
                "Importance of AI in today's world",
                "Applications in various industries"
            ]
        },
        {
            "title": "Types of AI",
            "content": [
                "Narrow AI",
                "General AI",
                "Superintelligent AI",
                "Examples of Narrow AI in use"
            ]
        },
        {
            "title": "AI and Machine Learning",
            "content": [
                "Overview of Machine Learning",
                "Difference between AI and ML",
                "Types of Machine Learning: Supervised, Unsupervised, Reinforcement",
                "Real-world applications of Machine Learning"
            ]
        },
        {
            "title": "Challenges in AI",
            "content": [
                "Ethical concerns",
                "Bias in AI models",
                "Data privacy issues",
                "Transparency and explainability"
            ]
        },
        {
            "title": "Future of AI",
            "content": [
                "Current trends in AI research and development",
                "Potential advancements in AI technologies",
                "Impact on job markets and industries",
                "Opportunities for future research"
            ]
        }
    ]

    template_path = 'template.pptx'  # Path to your PowerPoint template
    output_path = 'output.pptx'  # Path to save the generated presentation

    try:
        generate_presentation(slides_data, template_path, output_path)
        return {"success": "true", "ppt_path": "output.pptx"}
    except Exception as e:
        return {"success": "false", "error": str(e)}, 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=3000)